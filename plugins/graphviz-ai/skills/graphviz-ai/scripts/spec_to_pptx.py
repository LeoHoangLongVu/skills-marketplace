#!/usr/bin/env python3
"""
spec_to_pptx.py — convert a strict-diagram spec JSON into NATIVE PowerPoint shapes.

Same input as strict_svg_diagram.py (alignment keys resolved identically). Output is a
one-slide .pptx where every element is editable in PowerPoint:

  clusters  -> rounded rectangles + bold label text boxes (back layer)
  edges     -> one open freeform per edge (Manhattan polyline); hop bridges are
               approximated by short polyline arcs so crossings still read as jumps
  arrowheads-> small filled freeform triangles at the true endpoint
  nodes     -> rounded rectangles with centered multi-line text (first line bold)
  labels    -> white-filled text-box chips

Geometry is proportional to the SVG. Scale: pt-per-px = max(0.75, 12 / min_font), so the
NOMINAL PowerPoint font size never reads below 12pt. To make text LOOK bigger on the same
paper, raise the spec's min_font (e.g. 16px -> 12pt nominal at the natural 96 dpi paper
size) and re-layout — scaling the slide together with the font is visually a no-op. To drop the diagram
into a 16:9 deck, copy all shapes (Ctrl+A) and paste into yours - they arrive as one
scalable group.

  python3 spec_to_pptx.py spec.json out.pptx
"""
import json, math, sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import strict_svg_diagram as ssd

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

PX = 9525                      # EMU per px — set per-spec in convert()
PPP = 0.75                     # pt per px — ditto
OX = OY = 0                    # EMU offsets to center the diagram on the slide
MIN_PT = 12                    # nominal font floor — applies to --slide native only
SLIDE_16x9 = (12192000, 6858000)   # PowerPoint's default widescreen: 13.333 x 7.5 in
MARGIN_PT = 12                 # breathing room around the fitted diagram
EDGE = "1f3864"                # wire + arrowhead colour (matches the SVG)

def _rgb(h): return RGBColor.from_string(h.lstrip("#").lower())

def _X(v): return Emu(int(round(OX + v * PX)))   # position x
def _Y(v): return Emu(int(round(OY + v * PX)))   # position y
def _S(v): return Emu(int(round(v * PX)))        # size / extent

def _box(shapes, x, y, w, h, text, size, bold, color, align=PP_ALIGN.LEFT, fill=None):
    tb = shapes.add_textbox(_X(x), _Y(y), _S(w), _S(h))
    tf = tb.text_frame
    # kill spAutoFit: with wrap=none it re-centers the shrunk box, defeating algn
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = _rgb(color); f.name = "Arial"
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = _rgb(fill)
    return tb

def _freeform(shapes, pts, width_px, color, closed=False, fill=None):
    """Open (wire) or closed (arrowhead) freeform through px-coordinate points."""
    fb = shapes.build_freeform(_X(pts[0][0]), _Y(pts[0][1]), scale=1.0)
    fb.add_line_segments([(_X(x), _Y(y)) for x, y in pts[1:]], close=closed)
    sp = fb.convert_to_shape()
    sp.shadow.inherit = False
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        sp.line.fill.background()
    else:
        sp.fill.background()
        sp.line.color.rgb = _rgb(color); sp.line.width = Pt(width_px * PPP)
        # square caps keep right-angle joints crisp
        sp.line._get_or_add_ln().set("cap", "sq")
    return sp

def _wire_points(i, pts2, vsegs, r, hops_on, samples=8):
    """Edge polyline with each hop replaced by a sampled semicircular arc (bulge up)."""
    out = [pts2[0]]
    for P, Q in zip(pts2, pts2[1:]):
        if hops_on and P[1] == Q[1] and P[0] != Q[0]:
            y = P[1]; d = 1 if Q[0] > P[0] else -1
            xs = sorted({vx for (j, vx, ylo, yhi) in vsegs
                         if j != i and ylo < y < yhi and min(P[0], Q[0]) < vx < max(P[0], Q[0])},
                        reverse=(d < 0))
            for vx in xs:
                for t in range(samples + 1):
                    a = math.pi * t / samples
                    out.append((vx - d * r * math.cos(a), y - r * math.sin(a)))
        out.append(tuple(Q))
    return out

def _arrow_tri(e):
    (x0, y0), (x1, y1) = e["points"][-2], e["points"][-1]; t = 9
    if   x1 == x0 and y1 > y0: return [(x1, y1), (x1-4.5, y1-t), (x1+4.5, y1-t)]
    elif x1 == x0:             return [(x1, y1), (x1-4.5, y1+t), (x1+4.5, y1+t)]
    elif x1 > x0:              return [(x1, y1), (x1-t, y1-4.5), (x1-t, y1+4.5)]
    else:                      return [(x1, y1), (x1+t, y1-4.5), (x1+t, y1+4.5)]

def convert(spec_path, out_path, slide="16:9", fit=1.0):
    spec = ssd._resolve_align(json.load(open(spec_path)))
    W, H = spec.get("width", 2200), spec.get("height", 880)
    node_fs, chip_fs, cluster_fs, pitch = ssd._fonts(spec)
    global PX, PPP, OX, OY

    prs = Presentation()
    if slide == "native":
        # wall-chart: slide sized to the diagram, nominal font floor holds
        PPP = max(0.75, MIN_PT / node_fs)
        PX = 12700 * PPP
        OX = OY = 0
        prs.slide_width, prs.slide_height = Emu(int(W*PX)), Emu(int(H*PX))
    else:
        # DEFAULT: PowerPoint's own 13.333 x 7.5 in (16:9) slide — copy-paste into
        # another deck arrives at sane size. Diagram fits + centers; nominal font
        # follows from the fit (a 2560px-wide chart cannot read 12pt on one slide).
        # slide may also be: a .pptx path (match THAT deck's slide size exactly)
        # or "WxH" in inches (e.g. "10x7.5" for 4:3). fit<1 leaves margin for
        # titles: the pasted group occupies that fraction of the slide.
        if slide.endswith(".pptx"):
            t = Presentation(slide)
            sw, sh = int(t.slide_width), int(t.slide_height)
        elif slide not in ("16:9",) and "x" in slide:
            w_in, h_in = (float(v) for v in slide.lower().split("x"))
            sw, sh = int(w_in*914400), int(h_in*914400)
        else:
            sw, sh = SLIDE_16x9
        avail_w_pt = (sw/12700 - 2*MARGIN_PT) * fit
        avail_h_pt = (sh/12700 - 2*MARGIN_PT) * fit
        PPP = min(avail_w_pt / W, avail_h_pt / H)
        PX = 12700 * PPP
        OX = int((sw - W*PX) / 2)
        OY = int((sh - H*PX) / 2)
        prs.slide_width, prs.slide_height = Emu(sw), Emu(sh)
    shapes = prs.slides.add_slide(prs.slide_layouts[6]).shapes   # blank layout

    if spec.get("title"):
        _box(shapes, 0, 14, W, 34, spec["title"], max(22, node_fs+8)*PPP, False, "111111", PP_ALIGN.CENTER)

    for c in spec.get("clusters", []):
        w, h = c["x2"]-c["x1"], c["y2"]-c["y1"]
        sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _X(c["x1"]), _Y(c["y1"]), _S(w), _S(h))
        sp.adjustments[0] = 12 / min(w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(c.get("fill", "#f5f5f5"))
        sp.line.color.rgb = _rgb(c.get("stroke", "#999999")); sp.line.width = Pt(1.4*PPP)
        sp.shadow.inherit = False
        if c.get("label"):
            _box(shapes, c["x1"]+16, c["y1"]+8, w-32, cluster_fs+8, c["label"],
                 cluster_fs*PPP, True, c.get("labelcolor", "#333333"))

    drawn   = [ssd._drawn_points(e) for e in spec["edges"]]
    hops_on = spec.get("hops", True)
    vsegs   = ssd._vertical_segs(drawn) if hops_on else []
    hop_r   = spec.get("hop_radius", 6)
    for i, e in enumerate(spec["edges"]):
        _freeform(shapes, _wire_points(i, drawn[i], vsegs, hop_r, hops_on), 1.6, EDGE)
        _freeform(shapes, _arrow_tri(e), 0, EDGE, closed=True, fill=EDGE)

    for nid, n in spec["nodes"].items():
        x, y = n["cx"]-n["w"]/2, n["cy"]-n["h"]/2
        sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _X(x), _Y(y), _S(n["w"]), _S(n["h"]))
        sp.adjustments[0] = 9 / min(n["w"], n["h"])
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(n.get("fill", "#ffffff"))
        sp.line.color.rgb = _rgb(n.get("stroke", "#1f3864")); sp.line.width = Pt(1.6*PPP)
        sp.shadow.inherit = False
        tf = sp.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for j, line in enumerate(n.get("lines", [nid])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            # tighten line pitch to the SVG's rhythm (pitch px -> centipoints)
            pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
            spc = pPr.makeelement(qn("a:lnSpc"), {}); pts = spc.makeelement(qn("a:spcPts"), {"val": str(int(pitch*PPP*100))})
            spc.append(pts); pPr.insert(0, spc)
            r = p.add_run(); r.text = line
            f = r.font; f.name = "Arial"; f.color.rgb = _rgb("#111111")
            f.size = Pt(node_fs * PPP); f.bold = (j == 0)

    for e in spec["edges"]:
        if not e.get("label"): continue
        cx, cy = e.get("label_at", (0, 0))
        x1, y1, x2, y2 = ssd._chip_rect(e["label"], cx, cy, fs=chip_fs)
        _box(shapes, x1, y1, x2-x1, y2-y1, e["label"], chip_fs*PPP, False, "555555",
             PP_ALIGN.CENTER, fill="#ffffff")

    prs.save(out_path)
    sw_in, sh_in = prs.slide_width/914400, prs.slide_height/914400
    note = "" if slide == "native" else f" — fit {fit:.0%} of {sw_in:.2f}x{sh_in:.2f}in; --slide native for the wall-chart"
    print(f"wrote {out_path} ({sw_in:.2f}x{sh_in:.2f}in slide, body font {node_fs*PPP:.1f}pt nominal"
          f"{note}, {len(spec['nodes'])} nodes, {len(spec['edges'])} edges)")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.exit("usage: python3 spec_to_pptx.py spec.json out.pptx"
                 " [--slide 16:9|native|WxH|deck.pptx] [--fit 0.7]")
    slide = sys.argv[sys.argv.index("--slide")+1] if "--slide" in sys.argv else "16:9"
    fit = float(sys.argv[sys.argv.index("--fit")+1]) if "--fit" in sys.argv else 1.0
    convert(sys.argv[1], sys.argv[2], slide, fit)
